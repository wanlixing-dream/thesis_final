#!/usr/bin/env python3
"""生成毕业答辩PPT - InterMoBA-MTL"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── 配置 ───
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
IMG_DIR = os.path.join(BASE_DIR, "images")
OUTPUT_PATH = os.path.join(BASE_DIR, "答辩PPT_InterMoBA-MTL.pptx")

# 颜色方案（深蓝学术风格）
PRIMARY = RGBColor(0x1B, 0x3A, 0x6B)      # 深蓝
SECONDARY = RGBColor(0x2E, 0x86, 0xC1)    # 亮蓝
ACCENT = RGBColor(0xE7, 0x4C, 0x3C)       # 红色强调
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)


def set_slide_bg(slide, color):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title_text, left=Cm(0), top=Cm(0), width=Cm(25.4), height=Cm(2.2)):
    """添加标题栏"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    # 添加左侧缩进
    p.space_before = Pt(0)
    
    return shape


def add_text_box(slide, text, left, top, width, height, 
                 font_size=Pt(18), bold=False, color=BLACK, alignment=PP_ALIGN.LEFT):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def add_bullet_points(slide, points, left, top, width, height, 
                      font_size=Pt(16), color=BLACK, line_spacing=1.5):
    """添加要点列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # 支持加粗前缀（用 ** 包裹）
        if "**" in point:
            parts = point.split("**")
            for j, part in enumerate(parts):
                if part:
                    run = p.add_run()
                    run.text = part
                    run.font.size = font_size
                    run.font.color.rgb = color
                    run.font.bold = (j % 2 == 1)  # 奇数索引的部分加粗
        else:
            p.text = "• " + point
            p.font.size = font_size
            p.font.color.rgb = color
        
        p.space_after = Pt(8)
    
    return txBox


def add_image_safe(slide, img_name, left, top, width=None, height=None):
    """安全添加图片（文件不存在则跳过）"""
    img_path = os.path.join(IMG_DIR, img_name)
    if os.path.exists(img_path):
        kwargs = {"left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        slide.shapes.add_picture(img_path, **kwargs)
        return True
    else:
        print(f"  [警告] 图片不存在: {img_name}")
        return False


# ═══════════════════════════════════════════════════
# 开始生成PPT
# ═══════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = Cm(25.4)
prs.slide_height = Cm(19.05)

blank_layout = prs.slide_layouts[6]  # 空白布局

# ─── Slide 1: 封面 ───
print("生成 Slide 1: 封面")
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)

# 学校名称
add_text_box(slide, "深圳技术大学", Cm(1), Cm(1.5), Cm(23), Cm(1.5),
             font_size=Pt(20), color=PRIMARY, alignment=PP_ALIGN.CENTER)

# 论文题目
txBox = slide.shapes.add_textbox(Cm(2), Cm(4), Cm(21), Cm(4))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "InterMoBA-MTL：蛋白-配体多任务判别式\n预测与可解释能量建模的系统研究"
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = PRIMARY
p.alignment = PP_ALIGN.CENTER

# 作者信息
info_lines = [
    "答辩人：吴亮希",
    "指导教师：王鑫 助理教授",
    "学院：人工智能学院",
    "专业：计算机科学与技术",
]
txBox2 = slide.shapes.add_textbox(Cm(5), Cm(10), Cm(15), Cm(6))
tf2 = txBox2.text_frame
tf2.word_wrap = True
for i, line in enumerate(info_lines):
    if i == 0:
        p = tf2.paragraphs[0]
    else:
        p = tf2.add_paragraph()
    p.text = line
    p.font.size = Pt(18)
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(6)

# 日期
add_text_box(slide, "2026年5月", Cm(1), Cm(17), Cm(23), Cm(1.2),
             font_size=Pt(16), color=GRAY, alignment=PP_ALIGN.CENTER)

# 底部装饰线
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(18.3), Cm(25.4), Cm(0.75))
shape.fill.solid()
shape.fill.fore_color.rgb = PRIMARY
shape.line.fill.background()


# ─── Slide 2: 研究背景 ───
print("生成 Slide 2: 研究背景与动机")
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_title_bar(slide, "  一、研究背景与动机")

points = [
    "• **问题背景**：蛋白质-配体结合预测是计算机辅助药物设计的核心",
    "  → 一款新药平均耗资 20亿美元、耗时 10年以上",
    "",
    "• **现有局限**：深度学习方法通常围绕单一任务设计",
    "  → 亲和力预测 OR 姿态排序 OR 能量估计，缺乏统一框架",
    "",
    "• **核心挑战**：",
    "  → 不同任务数据源异构、量纲不同",
    "  → 共享主干上梯度方向冲突",
    "  → 简单 loss 加和效果不佳",
    "",
    "• **本文思路**：预训练+微调的两阶段多任务框架",
    '  → "先学结构，后学物理"',
]
txBox = slide.shapes.add_textbox(Cm(1.5), Cm(3), Cm(22), Cm(15))
tf = txBox.text_frame
tf.word_wrap = True
for i, pt in enumerate(points):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    
    if "**" in pt:
        parts = pt.split("**")
        for j, part in enumerate(parts):
            if part:
                run = p.add_run()
                run.text = part
                run.font.size = Pt(16)
                run.font.color.rgb = BLACK
                run.font.bold = (j % 2 == 1)
    else:
        p.text = pt
        p.font.size = Pt(16)
        p.font.color.rgb = BLACK if not pt.startswith("  →") else GRAY
    p.space_after = Pt(2)


# ─── Slide 3: 研究内容 ───
print("生成 Slide 3: 研究内容")
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_title_bar(slide, "  二、主要研究内容")

contents = [
    ("1", "Graph-Transformer 建模框架与预训练任务组织", 
     "基于 Interformer 共享主干，通过亲和力+姿态任务获得稳定表征"),
    ("2", "预训练—微调的两阶段训练策略", 
     "先冻结主干训练新任务头，再全参数联合优化"),
    ("3", "双数据源交替训练与任务掩蔽", 
     "Energy 与 Affinity+Pose 异构数据源交替采样"),
    ("4", "可解释能量建模与评测协议", 
     "四分量高斯能量分解 + RMSD/Pearson/AUROC 评测"),
]

y_start = Cm(3.5)
for i, (num, title, desc) in enumerate(contents):
    y = y_start + Cm(i * 3.5)
    
    # 编号圆圈
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(2), y, Cm(1.5), Cm(1.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = SECONDARY
    circle.line.fill.background()
    ctf = circle.text_frame
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = ctf.paragraphs[0]
    cp.text = num
    cp.font.size = Pt(18)
    cp.font.bold = True
    cp.font.color.rgb = WHITE
    cp.alignment = PP_ALIGN.CENTER
    
    # 标题
    add_text_box(slide, title, Cm(4.2), y - Cm(0.1), Cm(19), Cm(1.2),
                 font_size=Pt(17), bold=True, color=PRIMARY)
    # 描述
    add_text_box(slide, desc, Cm(4.2), y + Cm(1.1), Cm(19), Cm(1.2),
                 font_size=Pt(14), color=GRAY)


# ─── Slide 4: 系统架构 ───
print("生成 Slide 4: 系统架构")
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_title_bar(slide, "  三、InterMoBA-MTL 系统架构")

# 尝试插入系统架构图
if not add_image_safe(slide, "system_architecture.png", Cm(1.5), Cm(3), width=Cm(22)):
    # 如果图片不存在，用文字描述
    points = [
        "• 共享主干：12层 Graph-Transformer（Interformer）",
        "• 三路任务头：亲和力回归 / 姿态二分类 / 能量分解",
        "• 统一前向路径：一次 forward 输出所有任务结果",
        "• task_mask 控制损失分支激活",
        "• MoBA 路由：仅能量批次进入稀疏注意力路径",
    ]
    add_bullet_points(slide, points, Cm(2), Cm(4), Cm(21), Cm(12), font_size=Pt(16))


# ─── Slide 5: 两阶段训练 ───
print("生成 Slide 5: 两阶段训练策略")
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_title_bar(slide, "  四、两阶段训练策略")

# 左半部分：说明
txBox = slide.shapes.add_textbox(Cm(1), Cm(3), Cm(12.5), Cm(14))
tf = txBox.text_frame
tf.word_wrap = True

stage_info = [
    ("阶段一：预训练", PRIMARY),
    ("  • 亲和力回归 + 姿态选择", BLACK),
    ("  • 学习稳定的蛋白-配体交互表征", BLACK),
    ("  • Interformer 主干参数初始化", BLACK),
    ("", BLACK),
    ("阶段二：多任务微调", ACCENT),
    ("  • 引入能量监督（Vina 四分量）", BLACK),
    ("  • 前5 epoch 冻结主干", BLACK),
    ("  • 之后全参数联合优化", BLACK),
    ("  • 双数据源交替采样（p=0.3）", BLACK),
    ("", BLACK),
    ("关键效果：", SECONDARY),
    ("  • 预训练权重加载率 94.6%", BLACK),
    ("  • val_pose_sel: 0.50→0.801（epoch 0）", BLACK),
]

for i, (text, color) in enumerate(stage_info):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = color
    p.font.bold = text.endswith("：") or text.startswith("阶段")
    p.space_after = Pt(3)

# 右侧：权重调度图
add_image_safe(slide, "weight_scheduler.png", Cm(13.5), Cm(3.5), width=Cm(11))


# ─── Slide 6: 能量建模 ───
print("生成 Slide 6: 可解释能量建模")
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_title_bar(slide, "  五、可解释能量建模")

# 左侧：说明
txBox = slide.shapes.add_textbox(Cm(1), Cm(3), Cm(12), Cm(14))
tf = txBox.text_frame
tf.word_wrap = True

energy_info = [
    "四分量高斯混合能量模型：",
    "",
    "• vdW (attractive): 范德华吸引",
    "• vdW (repulsive): 范德华排斥",
    "• Hydrophobic: 疏水相互作用",
    "• H-bond: 氢键相互作用",
    "",
    "设计要点：",
    "• 原子对类型掩蔽机制",
    "  → 疏水项仅作用于非极性原子对",
    "  → 氢键项约束供体-受体匹配",
    "• 几何约束（碰撞检测+最小距离）",
    "• 提供原子对级可解释性",
]

for i, text in enumerate(energy_info):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = PRIMARY if text.endswith("：") else BLACK
    p.font.bold = text.endswith("：")
    p.space_after = Pt(3)

# 右侧图片
add_image_safe(slide, "gaussian_components.png", Cm(13), Cm(3), width=Cm(11.5))


# ─── Slide 7: 实验结果 - 亲和力 ───
print("生成 Slide 7: 实验结果 - 亲和力预测")
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_title_bar(slide, "  六、实验结果 — 亲和力预测")

# 结果表格（用文本模拟）
txBox = slide.shapes.add_textbox(Cm(1.5), Cm(3), Cm(13), Cm(9))
tf = txBox.text_frame
tf.word_wrap = True

table_lines = [
    "PDBbind 2020 time-split 测试集（341 targets）",
    "",
    "模型                    Pearson R    RMSE",
    "─────────────────────────────────",
    "GNINA                      0.495     1.74",
    "TankBind                   0.718     1.35",
    "Interformer (单模型)      0.692     1.37",
    "Interformer (4模型集成)   0.702     1.33",
    "InterMoBA-MTL (单模型)   0.737     1.27  ★",
    "",
    "★ 单模型超越4模型集成！",
    "  → 较最佳单模型提升 6.5%",
    "  → 较4模型集成提升 5.0%",
]

for i, text in enumerate(table_lines):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    if "InterMoBA" in text:
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = ACCENT
    elif text.startswith("★") or text.startswith("  →"):
        p.font.size = Pt(14)
        p.font.color.rgb = SECONDARY
        p.font.bold = text.startswith("★")
    else:
        p.font.size = Pt(13)
        p.font.color.rgb = BLACK

# 右侧散点图
add_image_safe(slide, "affinity_scatter_thesis.png", Cm(14), Cm(3), width=Cm(10.5))


# ─── Slide 8: 实验结果 - 姿态选择 & PoseBusters ───
print("生成 Slide 8: 姿态选择 & 独立验证")
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_title_bar(slide, "  七、姿态选择 & PoseBusters 独立验证")

# 左侧：姿态选择
txBox = slide.shapes.add_textbox(Cm(1), Cm(3), Cm(12), Cm(7))
tf = txBox.text_frame
tf.word_wrap = True

pose_info = [
    ("姿态选择（7161 poses, RMSD<2Å）", PRIMARY, True),
    ("", BLACK, False),
    ("模型              AUROC  Top-1%", BLACK, False),
    ("Interformer 集成   0.938  87.39%", BLACK, False),
    ("InterMoBA-MTL      0.924  88.56% ★", ACCENT, True),
    ("", BLACK, False),
    ("→ 单模型Top-1超越4模型集成", SECONDARY, False),
]

for i, (text, color, bold) in enumerate(pose_info):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_after = Pt(3)

# 左下：PoseBusters 说明
txBox2 = slide.shapes.add_textbox(Cm(1), Cm(11), Cm(12), Cm(6))
tf2 = txBox2.text_frame
tf2.word_wrap = True

pb_info = [
    ("PoseBusters v2 独立验证（n=308）", PRIMARY, True),
    ("", BLACK, False),
    ("• RMSD<2Å: 69.81%（所有方法最高）", BLACK, False),
    ("  → 超 Vina 9.8%, 超 DiffDock 31.8%", SECONDARY, False),
    ("• 证明分布外数据泛化能力", BLACK, False),
]

for i, (text, color, bold) in enumerate(pb_info):
    if i == 0:
        p = tf2.paragraphs[0]
    else:
        p = tf2.add_paragraph()
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_after = Pt(3)

# 右侧图片
add_image_safe(slide, "posebuster_benchmark_v2.png", Cm(13), Cm(3), width=Cm(11.5))


# ─── Slide 9: 可解释性 ───
print("生成 Slide 9: 可解释性分析")
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_title_bar(slide, "  八、可解释性 — 相互作用恢复分析")

# 左侧：恢复率数据
txBox = slide.shapes.add_textbox(Cm(1), Cm(3), Cm(12.5), Cm(14))
tf = txBox.text_frame
tf.word_wrap = True

interp_info = [
    ("PLIP 相互作用恢复率对比", PRIMARY, True, Pt(16)),
    ("", BLACK, False, Pt(12)),
    ("方法            氢键   疏水", BLACK, False, Pt(14)),
    ("──────────────────────", BLACK, False, Pt(14)),
    ("DeepDock         21.4%  ---", BLACK, False, Pt(14)),
    ("DiffDock         30.5%  ---", BLACK, False, Pt(14)),
    ("Interformer      50.2%  38.5%", BLACK, False, Pt(14)),
    ("InterMoBA-Energy 62.7%  50.1%", BLACK, False, Pt(14)),
    ("InterMoBA-MTL    79.3%  79.9% ★", ACCENT, True, Pt(14)),
    ("", BLACK, False, Pt(12)),
    ("关键发现：", SECONDARY, True, Pt(15)),
    ("• 能量监督显著提升局部交互敏感度", BLACK, False, Pt(14)),
    ("• 多任务联合训练有叠加效果", BLACK, False, Pt(14)),
    ("• 模型具备原子对级可解释输出能力", BLACK, False, Pt(14)),
]

for i, (text, color, bold, size) in enumerate(interp_info):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_after = Pt(3)

# 右侧：案例图
add_image_safe(slide, "6ggb_case.png", Cm(13.5), Cm(3), width=Cm(11))


# ─── Slide 10: 总结 ───
print("生成 Slide 10: 总结与贡献")
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)
add_title_bar(slide, "  九、总结与贡献")

txBox = slide.shapes.add_textbox(Cm(2), Cm(3.5), Cm(21), Cm(14))
tf = txBox.text_frame
tf.word_wrap = True

summary = [
    ("主要贡献", PRIMARY, True, Pt(20)),
    ("", BLACK, False, Pt(12)),
    ("❶  多任务统一架构", SECONDARY, True, Pt(17)),
    ("    三类异构任务在同一共享主干上联合训练，task_mask + MoBA路由实现软分离", BLACK, False, Pt(14)),
    ("", BLACK, False, Pt(8)),
    ("❷  两阶段训练策略 + 双数据源", SECONDARY, True, Pt(17)),
    ('    \u201c先学结构，后学物理\u201d，预训练权重正确加载使收敛加速', BLACK, False, Pt(14)),
    ("", BLACK, False, Pt(8)),
    ("❸  实验验证：单模型超越集成方案", SECONDARY, True, Pt(17)),
    ("    亲和力 R=0.737（+6.5%），Top-1 88.56%（超集成）", BLACK, False, Pt(14)),
    ("    PoseBusters 独立验证 RMSD<2Å 69.81%（最高）", BLACK, False, Pt(14)),
    ("", BLACK, False, Pt(8)),
    ("❹  可解释能量建模", SECONDARY, True, Pt(17)),
    ("    氢键/疏水恢复率 79%+，为药物优化提供原子级证据", BLACK, False, Pt(14)),
]

for i, (text, color, bold, size) in enumerate(summary):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_after = Pt(2)


# ─── Slide 11: 致谢 ───
print("生成 Slide 11: 致谢")
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, WHITE)

# 装饰条
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(25.4), Cm(3))
shape.fill.solid()
shape.fill.fore_color.rgb = PRIMARY
shape.line.fill.background()

# 主文字
add_text_box(slide, "感谢各位老师的指导与评审！", Cm(2), Cm(7), Cm(21), Cm(2.5),
             font_size=Pt(32), bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, "敬请批评指正", Cm(2), Cm(10.5), Cm(21), Cm(1.5),
             font_size=Pt(22), color=GRAY, alignment=PP_ALIGN.CENTER)

# 底部信息
add_text_box(slide, "答辩人：吴亮希  |  指导教师：王鑫  |  深圳技术大学 人工智能学院", 
             Cm(2), Cm(15), Cm(21), Cm(1.5),
             font_size=Pt(14), color=GRAY, alignment=PP_ALIGN.CENTER)

# 底部装饰
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(18.3), Cm(25.4), Cm(0.75))
shape.fill.solid()
shape.fill.fore_color.rgb = PRIMARY
shape.line.fill.background()


# ─── 保存 ───
prs.save(OUTPUT_PATH)
print(f"\n✅ PPT已生成: {OUTPUT_PATH}")
print(f"   共 {len(prs.slides)} 页幻灯片")
print(f"   建议答辩时间分配:")
print(f"   - 封面+背景: ~1分钟")
print(f"   - 研究内容+架构: ~1.5分钟")
print(f"   - 训练策略+能量建模: ~1分钟")
print(f"   - 实验结果: ~1分钟")
print(f"   - 总结: ~0.5分钟")
