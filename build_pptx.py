"""重新生成答辩PPT: InterMoBA-MTL"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── 路径 ───
BASE = '/home/wanlixing/桌面/intermoba_mtl/thesis_final'
IMG = os.path.join(BASE, 'images')
OUT = os.path.join(BASE, '答辩PPT_InterMoBA-MTL_v2.pptx')

# ─── 配色 ───
SZTU_BLUE = RGBColor(0x00, 0x3D, 0x87)  # 深技大蓝
ACCENT_BLUE = RGBColor(0x00, 0x72, 0xB2)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
GRAY = RGBColor(0x66, 0x66, 0x66)
ORANGE = RGBColor(0xE6, 0x9F, 0x00)

prs = Presentation()
prs.slide_width = Cm(25.4)
prs.slide_height = Cm(19.05)

SW = prs.slide_width
SH = prs.slide_height


def add_bg(slide, color=LIGHT_BG):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# 导航栏章节定义
NAV_SECTIONS = ['背景', '内容', '架构', '训练', '能量', '亲和力', '姿态', '可解释', '总结']
# 每页对应激活的导航栏 section index (第1页封面无导航，第2页起)
NAV_MAP = {2: -1, 3: 0, 4: 1, 5: 2, 6: 3, 7: 4, 8: 5, 9: 6, 10: 7, 11: 8}


def add_nav_bar(slide, page_num):
    """顶部导航栏：左侧章节导航 + 右侧校徽"""
    # 底色条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), SW, Cm(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = SZTU_BLUE
    bar.line.fill.background()

    active_idx = NAV_MAP.get(page_num, -1)
    n = len(NAV_SECTIONS)
    # 左侧 80% 给导航文字，右侧留给 logo
    nav_width = Cm(21.0)
    item_w = int(nav_width / n)

    for i, sec in enumerate(NAV_SECTIONS):
        left = int(item_w * i)
        txbox = slide.shapes.add_textbox(left, Cm(0.1), item_w, Cm(0.85))
        tf = txbox.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = sec
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(9)
        if i == active_idx:
            p.font.bold = True
            p.font.color.rgb = WHITE
        else:
            p.font.bold = False
            p.font.color.rgb = RGBColor(0xAA, 0xCC, 0xE8)

    # 右侧白色背景块（直角，上下顶满导航栏）+ 横版校名 logo
    bg_left = Cm(20.8)
    bg_w = SW - bg_left
    # 白色矩形顶满导航栏高度
    bg_rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, bg_left, Cm(0), bg_w, Cm(1.0))
    bg_rect.fill.solid()
    bg_rect.fill.fore_color.rgb = WHITE
    bg_rect.line.fill.background()
    # logo 居中放在白色块内
    logo_w = Cm(3.8)
    logo_h = Cm(0.75)
    logo_left = int(bg_left + (bg_w - logo_w) // 2)
    logo_top = int((Cm(1.0) - logo_h) // 2)
    slide.shapes.add_picture(
        os.path.join(IMG, 'school_title.png'),
        logo_left, logo_top, logo_w, logo_h)


def add_title_text(slide, title, top=Cm(1.5), font_size=Pt(28)):
    """添加标题"""
    txbox = slide.shapes.add_textbox(Cm(1.5), top, Cm(22), Cm(1.5))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = font_size
    p.font.bold = True
    p.font.color.rgb = SZTU_BLUE


def add_body_text(slide, lines, top=Cm(3.5), left=Cm(1.8), width=Cm(21), font_size=Pt(14)):
    """添加正文多行"""
    txbox = slide.shapes.add_textbox(left, top, width, Cm(13))
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.color.rgb = DARK
        p.space_after = Pt(6)
        # 加粗标记
        if line.startswith('★') or line.startswith('●'):
            p.font.bold = True
            p.font.color.rgb = ACCENT_BLUE


def add_image_centered(slide, img_path, top=Cm(4), max_w=Cm(20), max_h=Cm(13)):
    """居中添加图片，保持比例"""
    from PIL import Image
    im = Image.open(img_path)
    w_px, h_px = im.size
    aspect = w_px / h_px

    # 按最大宽高约束
    w = max_w
    h = int(w / aspect)
    if h > max_h:
        h = max_h
        w = int(h * aspect)

    left = (SW - w) // 2
    slide.shapes.add_picture(img_path, left, top, w, h)


def add_page_number(slide, num, total=12):
    """右下角页码"""
    txbox = slide.shapes.add_textbox(Cm(22.5), Cm(17.8), Cm(2.5), Cm(1))
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    p.text = f'{num}/{total}'
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT


# ================================================================
# 第1页：封面
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, SZTU_BLUE)

# Logo
logo_path = os.path.join(IMG, 'sztu.png')
slide.shapes.add_picture(logo_path, Cm(10.2), Cm(1.5), Cm(5), Cm(5.8))

# 标题
txbox = slide.shapes.add_textbox(Cm(1.5), Cm(8), Cm(22.4), Cm(3))
tf = txbox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = 'InterMoBA-MTL'
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = '蛋白-配体多任务判别式预测与可解释能量建模的系统研究'
p2.font.size = Pt(18)
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(8)

# 信息
txbox2 = slide.shapes.add_textbox(Cm(1.5), Cm(13), Cm(22.4), Cm(4))
tf2 = txbox2.text_frame
tf2.word_wrap = True
info_lines = [
    '答辩人：吴亮希',
    '指导教师：王鑫 助理教授',
    '学院：人工智能学院　专业：计算机科学与技术',
    '2026年5月'
]
for i, line in enumerate(info_lines):
    if i == 0:
        p = tf2.paragraphs[0]
    else:
        p = tf2.add_paragraph()
    p.text = line
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(4)

# ================================================================
# 第2页：目录
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_nav_bar(slide, 2)
add_title_text(slide, '目录')
add_page_number(slide, 2)

toc_items = [
    ('一', '研究背景与动机'),
    ('二', '主要研究内容'),
    ('三', 'InterMoBA-MTL 系统架构'),
    ('四', '两阶段训练策略'),
    ('五', '可解释能量建模'),
    ('六', '实验结果 — 亲和力预测'),
    ('七', '姿态选择 & PoseBusters 独立验证'),
    ('八', '可解释性 — 相互作用恢复分析'),
    ('九', '总结与贡献'),
]
top_y = Cm(3.8)
for idx, (num_label, title) in enumerate(toc_items):
    y = int(top_y + Cm(1.5) * idx)
    # 序号
    txbox_n = slide.shapes.add_textbox(Cm(2.5), y, Cm(2), Cm(1.2))
    tf_n = txbox_n.text_frame
    p_n = tf_n.paragraphs[0]
    p_n.text = num_label
    p_n.font.size = Pt(20)
    p_n.font.bold = True
    p_n.font.color.rgb = ACCENT_BLUE
    # 标题
    txbox_t = slide.shapes.add_textbox(Cm(4.8), y, Cm(18), Cm(1.2))
    tf_t = txbox_t.text_frame
    tf_t.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_t = tf_t.paragraphs[0]
    p_t.text = title
    p_t.font.size = Pt(15)
    p_t.font.color.rgb = DARK
    # 分隔线
    line_y = int(y + Cm(1.35))
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(2.5), line_y, Cm(20), Pt(0.5))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    line.line.fill.background()

# ================================================================
# 第3页：研究背景（左文右图）
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_nav_bar(slide, 3)
add_title_text(slide, '一、研究背景与动机')
add_page_number(slide, 3)

# ---- 左侧文字（占 55%）----
lines = [
    '● 问题背景',
    '   蛋白-配体结合预测是药物设计核心',
    '   新药平均耗资 20 亿美元、耗时 10+ 年',
    '',
    '● 现有局限',
    '   深度学习方法围绕单一任务设计',
    '   缺乏统一的多任务训练框架',
    '',
    '● 核心挑战',
    '   ① 数据源异构、量纲不同',
    '   ② 共享主干梯度方向冲突',
    '   ③ 简单 loss 加和效果不佳',
    '',
    '● 本文思路',
    '   "先学结构，后学物理"',
    '   → 预训练+微调两阶段框架',
]
add_body_text(slide, lines, top=Cm(3.2), left=Cm(1.2), width=Cm(12.5), font_size=Pt(12))

# ---- 右侧：单任务方法对比图 ----
# 三个任务色块 + 方法列表
card_colors = [
    (RGBColor(0x00, 0x72, 0xB2), RGBColor(0xE3, 0xEE, 0xF7), '亲和力预测', 'Pafnucy · PIGNet\nTankBind'),
    (RGBColor(0xE6, 0x9F, 0x00), RGBColor(0xFB, 0xF0, 0xDE), '姿态排序', 'DiffDock · GNINA\nEquiBind'),
    (RGBColor(0x00, 0x9E, 0x73), RGBColor(0xE8, 0xF2, 0xED), '能量估计', 'Vina 经验打分\nGOLD'),
]
card_left = Cm(14.5)
card_w = Cm(9.5)
card_h = Cm(3.0)
card_gap = Cm(0.5)
card_top_start = Cm(3.2)

for idx, (border_c, fill_c, task_name, methods) in enumerate(card_colors):
    y = int(card_top_start + (card_h + card_gap) * idx)
    # 卡片背景
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, card_left, y, card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = fill_c
    card.line.color.rgb = border_c
    card.line.width = Pt(1.5)
    # 左侧色条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, card_left, y, Cm(0.4), card_h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = border_c
    bar.line.fill.background()
    # 任务名
    txbox = slide.shapes.add_textbox(int(card_left + Cm(0.7)), y + Cm(0.3), Cm(8), Cm(1))
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    p.text = task_name
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = border_c
    # 方法列表
    txbox2 = slide.shapes.add_textbox(int(card_left + Cm(0.7)), y + Cm(1.3), Cm(8.5), Cm(1.5))
    tf2 = txbox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = methods
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY

# 虚线大括号区域 — 用文字标注 "各自为政"
brace_y = int(card_top_start + (card_h + card_gap) * 3 - Cm(0.3))
txbox_iso = slide.shapes.add_textbox(card_left, brace_y, card_w, Cm(1.2))
tf_iso = txbox_iso.text_frame
p_iso = tf_iso.paragraphs[0]
p_iso.text = '↑ 各自独立，缺乏统一框架'
p_iso.font.size = Pt(11)
p_iso.font.bold = True
p_iso.font.color.rgb = RGBColor(0xD5, 0x5E, 0x00)
p_iso.alignment = PP_ALIGN.CENTER

# 底部汇聚箭头 → InterMoBA-MTL
arrow_y = int(brace_y + Cm(1.2))
arrow = slide.shapes.add_shape(
    MSO_SHAPE.DOWN_ARROW, int(card_left + card_w // 2 - Cm(0.5)),
    arrow_y, Cm(1.0), Cm(1.0))
arrow.fill.solid()
arrow.fill.fore_color.rgb = SZTU_BLUE
arrow.line.fill.background()

# InterMoBA-MTL 统一框架标签
mtl_y = int(arrow_y + Cm(1.1))
mtl_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    int(card_left + Cm(1.0)), mtl_y, Cm(7.5), Cm(1.2))
mtl_box.fill.solid()
mtl_box.fill.fore_color.rgb = SZTU_BLUE
mtl_box.line.fill.background()
txbox_mtl = slide.shapes.add_textbox(
    int(card_left + Cm(1.0)), mtl_y, Cm(7.5), Cm(1.2))
tf_mtl = txbox_mtl.text_frame
tf_mtl.vertical_anchor = MSO_ANCHOR.MIDDLE
p_mtl = tf_mtl.paragraphs[0]
p_mtl.text = 'InterMoBA-MTL 多任务统一框架'
p_mtl.font.size = Pt(12)
p_mtl.font.bold = True
p_mtl.font.color.rgb = WHITE
p_mtl.alignment = PP_ALIGN.CENTER

# ================================================================
# 第3页：研究内容
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_nav_bar(slide, 4)
add_title_text(slide, '二、主要研究内容')
add_page_number(slide, 4)

lines = [
    '❶  Graph-Transformer 建模框架',
    '    基于 Interformer 共享主干，三类任务统一建模',
    '',
    '❷  预训练—微调的两阶段训练策略',
    '    先冻结主干训练新任务头，再全参数联合优化',
    '',
    '❸  双数据源交替训练与任务掩蔽',
    '    Energy 与 Affinity+Pose 异构数据源交替采样',
    '',
    '❹  可解释能量建模与评测协议',
    '    四分量高斯能量分解 + RMSD/Pearson/AUROC 评测',
]
add_body_text(slide, lines, font_size=Pt(14))

# ================================================================
# 第4页：系统架构
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_nav_bar(slide, 5)
add_title_text(slide, '三、InterMoBA-MTL 系统架构')
add_page_number(slide, 5)

arch_img = os.path.join(IMG, 'system_architecture.png')
add_image_centered(slide, arch_img, top=Cm(3.2), max_w=Cm(21), max_h=Cm(14.5))

# ================================================================
# 第5页：两阶段训练
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_nav_bar(slide, 6)
add_title_text(slide, '四、两阶段训练策略')
add_page_number(slide, 6)

lines = [
    '阶段一：预训练（Interformer 权重初始化）',
    '   • 亲和力回归 + 姿态选择',
    '   • 学习稳定的蛋白-配体交互表征',
    '',
    '阶段二：多任务微调',
    '   • 引入能量监督（Vina 四分量）',
    '   • 前 5 epoch 冻结编码器，只训练新任务头',
    '   • 之后全参数联合优化',
    '   • 双数据源交替采样（energy 比例 p=0.3）',
    '',
    '★ 关键效果：',
    '   • 预训练权重加载率 94.6%',
    '   • val_pose_sel: 0.50 → 0.801（epoch 0）',
    '   • 多任务增益：亲和力 R 从 0.692 提升到 0.737',
]
add_body_text(slide, lines, font_size=Pt(13))

# ================================================================
# 第6页：能量建模
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_nav_bar(slide, 7)
add_title_text(slide, '五、可解释能量建模')
add_page_number(slide, 7)

# 左侧文字
lines = [
    '四分量高斯混合距离头：',
    '  • vdW 吸引（全对开放）',
    '  • vdW 排斥（全对开放）',
    '  • 疏水作用（非极性门控）',
    '  • 氢键（供受体门控）',
    '',
    '设计要点：',
    '  • 原子对类型掩蔽',
    '  • 几何约束+碰撞检测',
    '  • 原子对级可解释输出',
]
add_body_text(slide, lines, top=Cm(3.5), left=Cm(1.5), width=Cm(10), font_size=Pt(12))

# 右侧图
gauss_img = os.path.join(IMG, 'gaussian_components.png')
from PIL import Image
im = Image.open(gauss_img)
w_px, h_px = im.size
aspect = w_px / h_px
img_w = Cm(13)
img_h = int(img_w / aspect)
slide.shapes.add_picture(gauss_img, Cm(11.5), Cm(3.5), img_w, img_h)

# ================================================================
# 第7页：亲和力结果
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_nav_bar(slide, 8)
add_title_text(slide, '六、实验结果 — 亲和力预测')
add_page_number(slide, 8)

# 表格文字
lines = [
    'PDBbind 2020 time-split 测试集（341 targets）',
    '',
    '  模型                         Pearson R    RMSE',
    '  ─────────────────────────────────────────',
    '  Interformer (单模型)            0.692      1.37',
    '  Interformer (4模型集成)         0.702      1.33',
    '★ InterMoBA-MTL (单模型)         0.737      1.27',
    '',
    '  → 单模型超越 4 模型集成，提升 5.0%',
]
add_body_text(slide, lines, top=Cm(3.0), font_size=Pt(13))

# scatter 图
scatter_img = os.path.join(IMG, 'affinity_scatter_thesis.png')
add_image_centered(slide, scatter_img, top=Cm(10.5), max_w=Cm(20), max_h=Cm(7.5))

# ================================================================
# 第8页：姿态选择 + PoseBusters
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_nav_bar(slide, 9)
add_title_text(slide, '七、姿态选择 & PoseBusters 独立验证')
add_page_number(slide, 9)

lines = [
    '姿态选择（7161 poses, RMSD<2Å 为正）',
    '  • Interformer 集成: AUROC 0.938, Top-1 87.39%',
    '★ InterMoBA-MTL:     AUROC 0.924, Top-1 88.56%',
    '  → 单模型 Top-1 超越 4 模型集成',
]
add_body_text(slide, lines, top=Cm(3.0), font_size=Pt(13))

# PoseBusters 图
pb_img = os.path.join(IMG, 'posebuster_benchmark_v2_nature.png')
add_image_centered(slide, pb_img, top=Cm(7.5), max_w=Cm(21), max_h=Cm(10))

# ================================================================
# 第9页：可解释性
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_nav_bar(slide, 10)
add_title_text(slide, '八、可解释性 — 相互作用恢复分析')
add_page_number(slide, 10)

# 上方文字
lines = [
    'PLIP 工具验证预测构象的相互作用恢复率',
    '',
    '★ InterMoBA-MTL: 氢键 79.3%, 疏水 79.9%',
    '   vs Interformer: 氢键 50.2%, 疏水 38.5%',
    '   → 能量监督 + 多任务联合训练显著提升局部交互敏感度',
]
add_body_text(slide, lines, top=Cm(3.0), font_size=Pt(13))

# 柱状图
recovery_img = os.path.join(IMG, 'conference_split/interaction_recovery_hist.png')
add_image_centered(slide, recovery_img, top=Cm(8.5), max_w=Cm(20), max_h=Cm(9))

# ================================================================
# 第10页：总结
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_nav_bar(slide, 11)
add_title_text(slide, '九、总结与贡献')
add_page_number(slide, 11)

lines = [
    '❶ 多任务统一架构',
    '   三类异构任务在同一共享主干上联合训练',
    '   task_mask 实现运行时任务路由',
    '',
    '❷ 两阶段训练策略 + 双数据源',
    '   "先学结构，后学物理"，预训练权重加速收敛',
    '',
    '❸ 实验验证：单模型全面超越集成方案',
    '   亲和力 R=0.737（+6.5%）',
    '   Top-1 88.56%（超集成）',
    '   PoseBusters RMSD<2Å 69.81%（所有方法最高）',
    '',
    '❹ 可解释能量建模',
    '   氢键/疏水恢复率 79%+',
    '   为药物优化提供原子级证据',
]
add_body_text(slide, lines, font_size=Pt(13))

# ================================================================
# 第11页：致谢
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SZTU_BLUE)

# Logo
slide.shapes.add_picture(logo_path, Cm(10.7), Cm(2), Cm(4), Cm(4.7))

txbox = slide.shapes.add_textbox(Cm(2), Cm(8), Cm(21.4), Cm(3))
tf = txbox.text_frame
p = tf.paragraphs[0]
p.text = '感谢各位老师的指导与评审！'
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = '敬请批评指正'
p2.font.size = Pt(20)
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(16)

txbox2 = slide.shapes.add_textbox(Cm(2), Cm(14), Cm(21.4), Cm(2))
tf2 = txbox2.text_frame
p3 = tf2.paragraphs[0]
p3.text = '答辩人：吴亮希  |  指导教师：王鑫  |  深圳技术大学 人工智能学院'
p3.font.size = Pt(12)
p3.font.color.rgb = WHITE
p3.alignment = PP_ALIGN.CENTER

# ─── 保存 ───
prs.save(OUT)
print(f'✓ PPT 已保存: {OUT}')
print(f'  共 {len(prs.slides)} 页')
